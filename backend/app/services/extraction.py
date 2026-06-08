"""Extraction du texte d'un document pour la recherche plein texte.

Stratégie en 2 étages :

1. **Extraction native** — pour les formats qui contiennent déjà du texte
   sélectionnable (PDF avec couche texte, Word, Excel, PowerPoint, ODT, TXT).
   Gratuit, instantané, qualité parfaite. C'est le cas dominant (70-90 % du
   volume en GED).

2. **OCR Tesseract** — pour les PDF scannés et les images. Utilise le binaire
   `tesseract` du système (à installer séparément, voir README).
   La langue par défaut est `fra` (français), modifiable via la variable
   d'env `OCR_LANGUES` (ex: `"fra+ara"` pour fr + arabe).

Le résultat est un texte UTF-8 brut, normalisé, prêt à être stocké dans
`documents.texte_ocr`. Le trigger Postgres se charge d'alimenter le
`recherche_fts` (tsvector) automatiquement.
"""

from __future__ import annotations

import io
import logging
import os
import re
from dataclasses import dataclass

logger = logging.getLogger(__name__)

# Limite de longueur du texte stocké en base. 1 million de caractères ≈
# 200 000 mots ≈ un livre entier. Au-delà, le texte n'a plus de valeur
# pratique pour la recherche et alourdit la DB inutilement.
LONGUEUR_MAX_TEXTE = 1_000_000

# Langues Tesseract — pack `fra` requis. Surcharge possible via env.
LANGUES_OCR_DEFAUT = os.getenv("OCR_LANGUES", "fra")

# Sur Windows, pytesseract a besoin du chemin absolu du binaire si Tesseract
# n'est pas dans le PATH système. On accepte une surcharge via env.
_TESSERACT_CMD = os.getenv("TESSERACT_CMD")


@dataclass(frozen=True, slots=True)
class ResultatExtraction:
    """Sortie d'une extraction de texte."""

    texte: str
    """Texte extrait, normalisé, tronqué à LONGUEUR_MAX_TEXTE caractères."""

    methode: str
    """Comment l'extraction a été obtenue : 'pdf_natif', 'pdf_ocr', 'docx',
    'xlsx', 'pptx', 'ocr_image', 'non_supporte'."""

    pages_traitees: int = 0
    """Nombre de pages OCR-isées (utile pour estimer le coût si on passe à
    une API payante un jour). 0 pour extraction native."""


class ExtractionError(RuntimeError):
    """Erreur d'extraction (fichier corrompu, format inattendu, etc.)."""


# ---------------------------------------------------------------------------
# Normalisation
# ---------------------------------------------------------------------------


_RE_ESPACES_MULTIPLES = re.compile(r"[ \t]+")
_RE_LIGNES_VIDES = re.compile(r"\n\s*\n+")


def _normaliser(texte: str) -> str:
    """Nettoie : espaces multiples, lignes vides multiples, troncature."""
    if not texte:
        return ""
    # Unifier les fins de ligne
    texte = texte.replace("\r\n", "\n").replace("\r", "\n")
    # Réduire les espaces et tabulations en série
    texte = _RE_ESPACES_MULTIPLES.sub(" ", texte)
    # Réduire les paquets de lignes vides à une seule ligne vide
    texte = _RE_LIGNES_VIDES.sub("\n\n", texte)
    texte = texte.strip()
    if len(texte) > LONGUEUR_MAX_TEXTE:
        texte = texte[:LONGUEUR_MAX_TEXTE]
    return texte


# ---------------------------------------------------------------------------
# Extractors par format
# ---------------------------------------------------------------------------


def _extraire_pdf(contenu: bytes) -> ResultatExtraction:
    """PDF : tente l'extraction native d'abord ; bascule sur OCR par page si
    le texte est vide (PDF scanné).

    Stratégie page-par-page pour gérer les PDF hybrides (certaines pages
    texte, d'autres scannées).
    """
    import pymupdf  # type: ignore[import-not-found]

    morceaux: list[str] = []
    pages_ocr = 0

    with pymupdf.open(stream=contenu, filetype="pdf") as doc:
        for page in doc:
            # 1) Tentative : extraction native
            texte_natif = page.get_text("text") or ""
            if len(texte_natif.strip()) >= 30:  # heuristique : page a du texte
                morceaux.append(texte_natif)
                continue

            # 2) Fallback OCR : on rend la page en image et on passe à Tesseract
            try:
                texte_ocr = _ocr_image_depuis_page(page)
                if texte_ocr:
                    morceaux.append(texte_ocr)
                    pages_ocr += 1
            except ExtractionError as exc:
                logger.warning(
                    "OCR page PDF échoué (page %s) : %s",
                    page.number,
                    exc,
                )

    texte = _normaliser("\n\n".join(morceaux))
    methode = "pdf_ocr" if pages_ocr > 0 else "pdf_natif"
    return ResultatExtraction(texte=texte, methode=methode, pages_traitees=pages_ocr)


def _ocr_image_depuis_page(page) -> str:  # type: ignore[no-untyped-def]
    """Rend une page PDF en bitmap puis appelle Tesseract."""
    # Pixmap à 200 DPI — bon compromis qualité OCR / vitesse
    pix = page.get_pixmap(dpi=200)
    return _ocr_octets_image(pix.tobytes("png"))


def _ocr_octets_image(contenu: bytes) -> str:
    """Lance Tesseract sur des octets d'image (PNG, JPG, etc.)."""
    try:
        import pytesseract  # type: ignore[import-not-found]
        from PIL import Image  # type: ignore[import-not-found]
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError(f"pytesseract / Pillow non installés : {exc}") from exc

    if _TESSERACT_CMD:
        pytesseract.pytesseract.tesseract_cmd = _TESSERACT_CMD

    try:
        with Image.open(io.BytesIO(contenu)) as img:
            # Mode "L" (niveaux de gris) améliore l'OCR sans dénaturer
            if img.mode not in ("L", "RGB"):
                img = img.convert("RGB")
            return pytesseract.image_to_string(img, lang=LANGUES_OCR_DEFAUT)
    except pytesseract.TesseractNotFoundError as exc:
        # Le binaire n'est pas installé : on retourne une chaîne vide plutôt
        # que de faire planter tout le pipeline. Le doc sera "ocr_echoue".
        raise ExtractionError(
            "Tesseract n'est pas installé ou n'est pas dans le PATH. "
            "Voir docs/install-tesseract.md."
        ) from exc
    except Exception as exc:  # pragma: no cover - garde-fou
        raise ExtractionError(f"Échec OCR Tesseract : {exc}") from exc


def _extraire_docx(contenu: bytes) -> ResultatExtraction:
    """Word .docx — extrait paragraphes + cellules de tableaux."""
    import docx  # type: ignore[import-not-found]

    doc = docx.Document(io.BytesIO(contenu))
    morceaux: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            morceaux.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            ligne = " | ".join(cell.text.strip() for cell in row.cells if cell.text)
            if ligne.strip():
                morceaux.append(ligne)
    return ResultatExtraction(
        texte=_normaliser("\n".join(morceaux)),
        methode="docx",
    )


def _extraire_xlsx(contenu: bytes) -> ResultatExtraction:
    """Excel .xlsx — concatène toutes les cellules non vides de toutes les
    feuilles. Pour un index FTS c'est suffisant ; on garde la structure
    visuelle simple (séparateur tab) pour ne pas générer de bruit."""
    import openpyxl  # type: ignore[import-not-found]

    wb = openpyxl.load_workbook(io.BytesIO(contenu), data_only=True, read_only=True)
    morceaux: list[str] = []
    for sheet in wb.worksheets:
        morceaux.append(f"# {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            valeurs = [
                str(v).strip() for v in row if v is not None and str(v).strip()
            ]
            if valeurs:
                morceaux.append("\t".join(valeurs))
    wb.close()
    return ResultatExtraction(
        texte=_normaliser("\n".join(morceaux)),
        methode="xlsx",
    )


def _extraire_pptx(contenu: bytes) -> ResultatExtraction:
    """PowerPoint .pptx — extrait les zones de texte de chaque slide."""
    from pptx import Presentation  # type: ignore[import-not-found]

    prs = Presentation(io.BytesIO(contenu))
    morceaux: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        morceaux.append(f"# Slide {i}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = "".join(run.text for run in para.runs)
                if txt.strip():
                    morceaux.append(txt)
    return ResultatExtraction(
        texte=_normaliser("\n".join(morceaux)),
        methode="pptx",
    )


def _extraire_image(contenu: bytes) -> ResultatExtraction:
    """Image (PNG, JPG, TIFF…) — OCR direct."""
    texte = _ocr_octets_image(contenu)
    return ResultatExtraction(
        texte=_normaliser(texte),
        methode="ocr_image",
        pages_traitees=1,
    )


def _extraire_texte_brut(contenu: bytes) -> ResultatExtraction:
    """text/plain — décodage UTF-8 avec fallback latin-1."""
    try:
        texte = contenu.decode("utf-8")
    except UnicodeDecodeError:
        texte = contenu.decode("latin-1", errors="ignore")
    return ResultatExtraction(texte=_normaliser(texte), methode="texte_brut")


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

# Mapping MIME → extracteur. Si le MIME n'est pas listé, on retourne un
# ResultatExtraction vide avec méthode 'non_supporte' (le doc reste utilisable,
# juste pas indexé en plein texte).
_EXTRACTEURS = {
    "application/pdf": _extraire_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extraire_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extraire_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extraire_pptx,
    "text/plain": _extraire_texte_brut,
    "text/markdown": _extraire_texte_brut,
    "text/csv": _extraire_texte_brut,
}


def extraire_texte(contenu: bytes, mime: str) -> ResultatExtraction:
    """Point d'entrée principal — dispatche selon le type MIME.

    Args:
        contenu: octets en clair du fichier (déchiffré au préalable).
        mime: type MIME exact (résultat de la détection serveur).

    Returns:
        Un ResultatExtraction toujours valide. Texte vide si le format
        n'est pas supporté ou si l'extraction échoue silencieusement.

    Raises:
        ExtractionError: en cas de fichier corrompu non récupérable. Les
            erreurs OCR (Tesseract absent, etc.) sont remontées aussi —
            le caller décide quoi en faire (statut "ocr_echoue" + log).
    """
    extracteur = _EXTRACTEURS.get(mime)

    if extracteur is not None:
        return extracteur(contenu)

    # Toutes les images passent par OCR
    if mime.startswith("image/"):
        return _extraire_image(contenu)

    # Format inconnu — on n'indexe pas mais on ne plante pas
    return ResultatExtraction(texte="", methode="non_supporte")
